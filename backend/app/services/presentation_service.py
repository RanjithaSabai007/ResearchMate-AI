import json
import re
import os
from datetime import datetime
import ollama
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def generate_presentation_outline(draft_title: str, draft_content: str) -> list:
    """
    Query local Ollama model to summarize draft content into structured slides.
    Returns a list of dictionaries with 'title' and 'bullets'.
    """
    prompt = f"""You are a university academic advisor and presentation slide designer.
Your task is to analyze the student's thesis draft and generate a structured presentation outline.
You MUST extract details ONLY from the provided draft content. Do not invent any facts, claims, methodologies, datasets, algorithms, results, future work, or conclusions.

Thesis Title: {draft_title}
Thesis Content:
{draft_content}

Generate a structured outline of the slides representing this thesis.
Omit slides for categories (like Research Gap, Implementation, Results, or Future Work) if they are not explicitly discussed or present in the thesis draft. Do not fabricate content.

You MUST respond with a single, valid JSON array containing slide objects matching this structure EXACTLY:
[
  {{
    "title": "Problem Statement & Motivation",
    "bullets": [
      "Key point summarizing the problem...",
      "Why this research is necessary..."
    ]
  }},
  {{
    "title": "Research Objectives",
    "bullets": [
      "Objective 1...",
      "Objective 2..."
    ]
  }}
]

Strict Constraints:
1. Output ONLY a valid raw JSON array.
2. Do NOT wrap the JSON in markdown blocks like ```json or ```.
3. Do NOT include any intro, outro, explanations, or commentary.
4. Each content slide must have 3 to 7 bullet points. Keep bullets short and presentation-friendly.
"""

    try:
        print("[INFO] Contacting local Ollama model (phi3:mini) for slide outline...")
        response = ollama.chat(
            model="phi3:mini",
            messages=[
                {
                    "role": "user",
                    "content": prompt
                }
            ]
        )
        content = response["message"]["content"].strip()
        
        # Clean markdown wrappers if present
        if content.startswith("```"):
            content = re.sub(r"^```(?:json)?\n", "", content)
            content = re.sub(r"\n```$", "", content)
            content = content.strip()
            
        outline = json.loads(content)
        if isinstance(outline, list) and len(outline) > 0:
            print(f"[INFO] Successfully parsed slide outline with {len(outline)} slides.")
            return outline
        else:
            raise ValueError("Ollama outline is empty or not a list")
    except Exception as e:
        print(f"[WARNING] Ollama presentation outline generation failed: {e}. Using fallback regex generator.")
        return fallback_outline_generator(draft_content)

def fallback_outline_generator(content: str) -> list:
    """
    Regex-based fallback slide parser that extracts sections directly from the draft.
    Ensure we never fail to generate a presentation.
    """
    # Clean HTML tags if present (e.g. from rich text editor)
    clean_text = re.sub(r'<[^>]*>', '\n', content)
    paras = [p.strip() for p in re.split(r'\n+', clean_text) if p.strip()]
    
    outline = []
    motivation = []
    objectives = []
    methodology = []
    results = []
    conclusions = []
    
    for p in paras:
        # Ignore extremely short strings
        if len(p) < 15:
            continue
        p_clean = p[:100] + ("..." if len(p) > 100 else "")
        p_lower = p.lower()
        
        if any(x in p_lower for x in ["motivation", "problem statement", "introduction", "background", "need"]):
            if len(motivation) < 5:
                motivation.append(p_clean)
        elif any(x in p_lower for x in ["objective", "goal", "aim", "target"]):
            if len(objectives) < 5:
                objectives.append(p_clean)
        elif any(x in p_lower for x in ["method", "architecture", "design", "algorithm", "dataset", "framework"]):
            if len(methodology) < 5:
                methodology.append(p_clean)
        elif any(x in p_lower for x in ["result", "experiment", "accuracy", "performance", "metric", "table"]):
            if len(results) < 5:
                results.append(p_clean)
        elif any(x in p_lower for x in ["conclusion", "conclude", "future work", "recommendation"]):
            if len(conclusions) < 5:
                conclusions.append(p_clean)

    # Build slides based on extracted bullet lists
    if motivation:
        outline.append({"title": "Motivation & Problem Statement", "bullets": motivation})
    else:
        outline.append({
            "title": "Motivation & Problem Statement",
            "bullets": [
                "Detailed overview of the research domain goals.",
                "Outlines primary challenges solved by the thesis draft."
            ]
        })
        
    if objectives:
        outline.append({"title": "Research Objectives", "bullets": objectives})
        
    if methodology:
        outline.append({"title": "Proposed Methodology", "bullets": methodology})
        
    if results:
        outline.append({"title": "Experimental Results & Validation", "bullets": results})
        
    if conclusions:
        outline.append({"title": "Research Conclusions", "bullets": conclusions})
        
    return outline

def create_pptx_from_outline(draft_title: str, slides_outline: list) -> str:
    """
    Creates a PPTX file using python-pptx from the structured outline list.
    Returns the absolute path to the generated presentation.
    """
    prs = Presentation()
    
    # 16:9 aspect ratio standard
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Theme color definitions
    PRIMARY_BG = RGBColor(15, 23, 42)     # Slate 900
    ACCENT_TITLE = RGBColor(99, 102, 241)  # Indigo 500
    BODY_TEXT = RGBColor(71, 85, 105)      # Slate 600
    WHITE = RGBColor(255, 255, 255)
    DIVIDER_LINE = RGBColor(226, 232, 240) # Slate 200
    
    # --- SLIDE 1: Title Slide ---
    title_slide = prs.slides.add_slide(blank_layout)
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = PRIMARY_BG
    
    # Centered Title box
    title_box = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = draft_title or "Thesis Presentation Defense"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.text = "\nAcademic Presentation generated by ResearchMate AI"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
    p_sub.alignment = PP_ALIGN.CENTER
    
    # --- CONTENT SLIDES ---
    for slide_data in slides_outline:
        title_text = slide_data.get("title", "Untitled Slide")
        bullets = slide_data.get("bullets", [])
        if not bullets:
            continue
            
        slide = prs.slides.add_slide(blank_layout)
        
        # Header Text Box
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.0))
        tf_header = header_box.text_frame
        tf_header.word_wrap = True
        p_header = tf_header.paragraphs[0]
        p_header.text = title_text
        p_header.font.size = Pt(32)
        p_header.font.bold = True
        p_header.font.color.rgb = ACCENT_TITLE
        
        # Horizontal divider line
        divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.03))
        divider.fill.solid()
        divider.fill.fore_color.rgb = DIVIDER_LINE
        divider.line.color.rgb = DIVIDER_LINE
        
        # Content Text Box
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        for idx, bullet in enumerate(bullets):
            p_bullet = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
            p_bullet.text = f"•  {bullet}"
            p_bullet.font.size = Pt(20)
            p_bullet.font.color.rgb = BODY_TEXT
            p_bullet.space_after = Pt(16)
            p_bullet.level = 0
            
    # --- FINAL SLIDE: Thank You ---
    thank_slide = prs.slides.add_slide(blank_layout)
    thank_slide.background.fill.solid()
    thank_slide.background.fill.fore_color.rgb = PRIMARY_BG
    
    thank_box = thank_slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(2.5))
    tf_thank = thank_box.text_frame
    tf_thank.word_wrap = True
    p_thank = tf_thank.paragraphs[0]
    p_thank.text = "Thank You!"
    p_thank.font.size = Pt(50)
    p_thank.font.bold = True
    p_thank.font.color.rgb = WHITE
    p_thank.alignment = PP_ALIGN.CENTER
    
    p_ques = tf_thank.add_paragraph()
    p_ques.text = "\nQuestions & Discussions"
    p_ques.font.size = Pt(22)
    p_ques.font.color.rgb = ACCENT_TITLE
    p_ques.alignment = PP_ALIGN.CENTER
    
    # Save the file
    os.makedirs("temp_presentations", exist_ok=True)
    filename = f"temp_presentations/presentation_{int(datetime.now().timestamp())}.pptx"
    prs.save(filename)
    return filename
